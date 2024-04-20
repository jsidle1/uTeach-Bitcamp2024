import os
import openai
import json


def run_gpt_model():
    openai.api_key = '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'

    completion = openai.chat.completions.create(
        model = "ft:gpt-3.5-turbo-0125:personal::9G6A0hni",
        messages = [
            {"role":"user","content": "I want to learn doppler effect, Can you create a presentation"}
        ],
        n = 1,
        stop = None,
        temperature = 0.7,
    )
    return completion.choices[0].message.content 
chatGPT_output = run_gpt_model()



#chatGPT_output = completion.choices[0].message.content
count = 0
while count < 10:
    try: 
        chatGPT_data_json = json.loads(chatGPT_output)
        break
    except ValueError:
        print("error has occur, rerunning")
        chatGPT_output = run_gpt_model()
if count >= 10:
    print("ERROR, ENDED")



#Getting pptx imported
from pptx import Presentation
#craete the introduction page
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = chatGPT_data_json['title']
subtitle.text = "Name: "

len_of_slides = len(chatGPT_data_json['slides'])


#Generate all the slides
for x in range(len_of_slides):
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    #Title of the slide
    title_shape.text = chatGPT_data_json['slides'][x]['title']

    #getting the bullet point
    tf = body_shape.text_frame
    tf.text = chatGPT_data_json['slides'][x]['content']
prs.save('text.pptx')