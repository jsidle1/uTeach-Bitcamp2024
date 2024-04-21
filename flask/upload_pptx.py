from flask import Flask, jsonify, request, send_file
from flask_cors import CORS
from werkzeug.utils import secure_filename
from scoreSpeech import scoreText

from pptx import Presentation

import subprocess
import os
import openai
import json

app = Flask(__name__)
cors = CORS(app)

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        print("here1")
        return jsonify({'error': 'No file part'}), 400
    file = request.files['file']
    if file.filename == '':
        print("here2")
        return jsonify({'error': 'No selected file'}), 400
    if file and allowed_file(file.filename):
        filename = secure_filename(file.filename)
        file.save(os.path.join(filename))
        #file.save(os.path.join('/path/to/save', filename))
        #return jsonify({'message': 'File uploaded successfully', 'filename': filename}), 200
            # Open the presentation
        ppt = Presentation(filename)
        
        # Extract text from each slide in the presentation
        pptText = []
        for slide in ppt.slides:
            slide_text = ""
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        slide_text += run.text
            pptText.append(slide_text)
        
        # Score each extracted text
        scores = []
        for slide_text in pptText:
            scores.append([(scoreText(slide_text, "correctness")), (scoreText(slide_text, "relevance")), (scoreText(slide_text, "clarity")), (scoreText(slide_text, "grammar"))])
        
        return jsonify({'message': scores})

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in {'pptx'}

if __name__ == '__main__':
    app.run(debug=True)