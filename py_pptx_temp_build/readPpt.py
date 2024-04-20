from scoreSpeech import scoreText
from pptx import Presentation 
  
# creating an object 
ppt = Presentation("py_pptx_temp_build/testPpt.pptx") 
  
# write text from powerpoint 
# file into .txt file 

pptText = []
for slide in ppt.slides: 
    slide_text = "" 
    for shape in slide.shapes:  
        if not shape.has_text_frame:  
            continue 
        for paragraph in shape.text_frame.paragraphs:  
            for run in paragraph.runs:  
                slide_text += run.text
    pptText.append(slide_text)

for slide_text in pptText:
    print(slide_text)
    score = scoreText(slide_text, "correctness")
    print("Score: " + score[0], "\nImprove by: " + score[1] + "\n\n") 